import tkinter as tk
from tkinter import ttk, scrolledtext, messagebox, filedialog # 导入 filedialog
import requests
from bs4 import BeautifulSoup
import datetime
import os
from docx import Document
from docx.shared import Cm as word_Cm
from pptx import Presentation
from pptx.util import Cm as ppt_Cm
from PIL import Image # Pillow库，用于读取图片尺寸
import threading
import queue

# --- 常量定义 (来自原始 weixin.py) ---
USER_AGENT = 'Mozilla/5.0 (iPhone; CPU iPhone OS 11_0 like Mac OS X) AppleWebKit/604.1.38 (KHTML, like Gecko) Version/11.0 Mobile/15A372 Safari/604.1'
# BASE_DESKTOP_FOLDER_NAME = "weixin_images_gui" # 不再固定到桌面，此常量可以移除或修改用途
DEFAULT_IMAGE_EXTENSION = "jpg"

# --- 核心逻辑函数 (从 weixin.py 修改而来) ---

def log_status(status_queue, message):
    """将状态消息放入队列，以便UI线程安全地更新文本框"""
    print(message) # 也在控制台打印一份，方便调试
    status_queue.put(message)

def create_timestamped_folder(status_queue, base_save_path: str): # 接收基础保存路径
    """
    在指定的基础路径下创建一个带时间戳的子文件夹。
    返回创建的子文件夹的绝对路径。
    """
    curr_time = datetime.datetime.now()
    timestamp_str = datetime.datetime.strftime(curr_time, '%Y%m%d%H%M%S')

    if not base_save_path: # 如果没有提供基础路径，可以设置一个默认值或报错
        log_status(status_queue, "错误：未指定保存路径。")
        # 或者，可以退回到桌面创建：
        # base_save_path = os.path.join(os.path.expanduser("~"), "Desktop", "weixin_images_gui_default")
        # log_status(status_queue, f"警告：未指定保存路径，将使用默认路径: {base_save_path}")
        return None


    main_folder_path = base_save_path # 用户选择的路径作为主文件夹

    try:
        if not os.path.exists(main_folder_path):
            # 通常用户选择的文件夹应该存在，如果不存在，根据需求看是否创建或报错
            # os.makedirs(main_folder_path) # 如果需要创建基础路径
            log_status(status_queue, f"警告：选择的基础保存路径 '{main_folder_path}' 不存在，请确保路径有效。")
            # return None # 如果路径必须存在则返回None

        session_folder_path = os.path.join(main_folder_path, timestamp_str)
        if not os.path.exists(session_folder_path):
            os.makedirs(session_folder_path)
        log_status(status_queue, f"文件将保存在: {session_folder_path}")
        return session_folder_path
    except Exception as e:
        log_status(status_queue, f"错误：创建文件夹失败 - {e}")
        return None


def download_images_from_url(url: str, save_folder: str, status_queue) -> list[str]:
    """
    从给定的微信公众号URL下载图片到指定的文件夹。
    返回成功下载的图片文件的完整路径列表。
    """
    headers = {'user-agent': USER_AGENT}
    downloaded_image_paths = []

    log_status(status_queue, f"开始从URL下载图片: {url}")
    try:
        response = requests.get(url=url, headers=headers, timeout=30) # 增加超时
        response.raise_for_status()
        html_content = response.content.decode('utf-8', errors='ignore')
    except requests.exceptions.RequestException as e:
        log_status(status_queue, f"错误：请求URL失败 - {url}, {e}")
        return downloaded_image_paths

    soup = BeautifulSoup(html_content, 'lxml')
    image_tags = soup.select('img') # 主要选择img标签

    if not image_tags:
        log_status(status_queue, "未在页面中找到 <img> 标签。")
        # 尝试查找可能的背景图片或其他形式的图片，这部分比较复杂，暂时简化
        # for style_tag in soup.find_all('style'):
        #    pass # 更复杂的CSS背景图提取逻辑

    image_counter = 0
    total_images_found = len(image_tags)
    log_status(status_queue, f"检测到 {total_images_found} 个图片标签。开始下载...")

    for i, img_tag in enumerate(image_tags):
        img_data_src = img_tag.get("data-src") or img_tag.get("src") # 兼容data-src和src
        if not img_data_src:
            # log_status(status_queue, f"跳过一个没有data-src或src属性的图片标签 ({i+1}/{total_images_found})")
            continue

        # 确保URL是完整的
        if img_data_src.startswith('//'):
            img_data_src = 'http:' + img_data_src # 或者 'https:'，根据实际情况
        elif not img_data_src.startswith(('http://', 'https://')):
            # log_status(status_queue, f"跳过无效的图片URL: {img_data_src}")
            continue


        img_extension = img_tag.get("data-type", DEFAULT_IMAGE_EXTENSION).split('/')[-1] # 如 image/jpeg -> jpeg
        if not img_extension or len(img_extension) > 5 : # 简单过滤无效扩展名
            img_extension = DEFAULT_IMAGE_EXTENSION
            # 尝试从URL中获取扩展名
            try:
                parsed_ext = os.path.splitext(img_data_src.split('?')[0])[-1].lstrip('.')
                if parsed_ext and len(parsed_ext) <= 4:
                    img_extension = parsed_ext
            except:
                pass

        img_filename = f"{image_counter}.{img_extension}"
        img_full_path = os.path.join(save_folder, img_filename)

        try:
            log_status(status_queue, f"下载中 ({image_counter + 1}/{total_images_found}): {img_data_src[:70]}...")
            img_response = requests.get(url=img_data_src, headers=headers, timeout=20)
            img_response.raise_for_status()
            with open(img_full_path, 'wb') as f:
                f.write(img_response.content)
            downloaded_image_paths.append(img_full_path)
            image_counter += 1
        except requests.exceptions.RequestException as e:
            log_status(status_queue, f"警告：下载图片失败 - {img_data_src[:70]}..., {e}")
        except IOError as e:
            log_status(status_queue, f"警告：保存图片失败 - {img_full_path}, {e}")
        except Exception as e:
            log_status(status_queue, f"警告：处理图片时发生未知错误 - {img_data_src[:70]}..., {e}")

    log_status(status_queue, f"图片下载完成，此次共成功保存 {image_counter} 张到文件夹: {save_folder}")
    return downloaded_image_paths


def generate_word_document(file_name_prefix: str, image_paths: list[str], save_folder: str, status_queue):
    if not image_paths:
        log_status(status_queue, "没有图片可用于生成Word文档。")
        return None

    log_status(status_queue, "开始生成Word文档...")
    doc = Document()

    for section in doc.sections:
        section.left_margin = word_Cm(0.5) # 留一些边距
        section.right_margin = word_Cm(0.5)
        section.top_margin = word_Cm(0.5)
        section.bottom_margin = word_Cm(0.5)

    page_width_cm = doc.sections[0].page_width.cm - (section.left_margin.cm + section.right_margin.cm)

    for i, img_path in enumerate(image_paths):
        try:
            log_status(status_queue, f"添加图片到Word ({i+1}/{len(image_paths)}): {os.path.basename(img_path)}")
            # 尝试获取图片原始尺寸以保持宽高比
            with Image.open(img_path) as img:
                width_px, height_px = img.size

            aspect_ratio = height_px / width_px
            # display_height_cm = page_width_cm * aspect_ratio # 变量未使用，移除

            doc.add_picture(img_path, width=word_Cm(page_width_cm)) # 高度会自动按比例调整
            # 如果希望每张图片后分页：
            # if i < len(image_paths) - 1:
            #    doc.add_page_break()
        except Exception as e:
            log_status(status_queue, f"警告：无法将图片添加到Word - {os.path.basename(img_path)}, {e}")

    output_filename = f"{file_name_prefix}.docx"
    output_full_path = os.path.join(save_folder, output_filename)
    try:
        doc.save(output_full_path)
        log_status(status_queue, f"Word文档已成功保存到: {output_full_path}")
        return output_full_path
    except Exception as e:
        log_status(status_queue, f"错误：保存Word文档失败 - {output_full_path}, {e}")
        return None


def generate_ppt_presentation(file_name_prefix: str, image_paths: list[str], save_folder: str, status_queue):
    if not image_paths:
        log_status(status_queue, "没有图片可用于生成PPT。")
        return None

    log_status(status_queue, "开始生成PPT演示文稿...")
    prs = Presentation()
    # 使用16:9的幻灯片尺寸，更常见
    prs.slide_width = ppt_Cm(33.867) # 16:9 width
    prs.slide_height = ppt_Cm(19.05) # 16:9 height

    slide_width_cm = prs.slide_width.cm
    slide_height_cm = prs.slide_height.cm
    slide_aspect_ratio = slide_width_cm / slide_height_cm

    for i, img_path in enumerate(image_paths):
        try:
            log_status(status_queue, f"添加图片到PPT ({i+1}/{len(image_paths)}): {os.path.basename(img_path)}")
            slide_layout = prs.slide_layouts[5] # 空白布局
            slide = prs.slides.add_slide(slide_layout)

            with Image.open(img_path) as img:
                img_width_px, img_height_px = img.size

            if img_height_px == 0 or img_width_px == 0: # 避免除以零
                log_status(status_queue, f"警告: 图片尺寸为零，跳过 {os.path.basename(img_path)}")
                continue
            img_aspect_ratio = img_width_px / img_height_px

            if img_aspect_ratio > slide_aspect_ratio:
                pic_display_width_cm = slide_width_cm
                pic_display_height_cm = pic_display_width_cm / img_aspect_ratio
            else:
                pic_display_height_cm = slide_height_cm
                pic_display_width_cm = pic_display_height_cm * img_aspect_ratio

            left_cm = (slide_width_cm - pic_display_width_cm) / 2
            top_cm = (slide_height_cm - pic_display_height_cm) / 2

            slide.shapes.add_picture(
                img_path,
                ppt_Cm(left_cm),
                ppt_Cm(top_cm),
                width=ppt_Cm(pic_display_width_cm),
                height=ppt_Cm(pic_display_height_cm)
            )
        except Exception as e:
            log_status(status_queue, f"警告：无法将图片添加到PPT - {os.path.basename(img_path)}, {e}")

    output_filename = f"{file_name_prefix}.pptx"
    output_full_path = os.path.join(save_folder, output_filename)
    try:
        prs.save(output_full_path)
        log_status(status_queue, f"PPT演示文稿已成功保存到: {output_full_path}")
        return output_full_path
    except Exception as e:
        log_status(status_queue, f"错误：保存PPT失败 - {output_full_path}, {e}")
        return None

# --- UI相关的类和函数 ---

class WeixinToolApp:
    def __init__(self, root):
        self.root = root
        self.root.title("微信公众号文章处理工具 v1.2") # 版本号更新
        self.root.geometry("700x600") # 稍微增加高度以容纳新控件

        # 用于线程通信的状态队列
        self.status_queue = queue.Queue()
        self.selected_save_path = "" # 用于存储用户选择的保存路径

        # --- 界面元素 ---
        # URL输入
        ttk.Label(root, text="微信文章 URL:").grid(row=0, column=0, padx=10, pady=10, sticky="w")
        self.url_entry = ttk.Entry(root, width=70)
        self.url_entry.grid(row=0, column=1, columnspan=2, padx=10, pady=10, sticky="ew") # columnspan 2

        # 文档名称前缀输入
        ttk.Label(root, text="文档名称前缀:").grid(row=1, column=0, padx=10, pady=5, sticky="w")
        self.prefix_entry = ttk.Entry(root, width=70)
        self.prefix_entry.grid(row=1, column=1, columnspan=2, padx=10, pady=5, sticky="ew") # columnspan 2
        self.prefix_entry.insert(0, "微信文章") # 默认值

        # 保存位置选择
        ttk.Label(root, text="选择保存文件夹:").grid(row=2, column=0, padx=10, pady=5, sticky="w")
        self.save_path_entry = ttk.Entry(root, width=55, state="readonly") # 显示路径的文本框
        self.save_path_entry.grid(row=2, column=1, padx=(0,5), pady=5, sticky="ew")
        self.browse_button = ttk.Button(root, text="浏览...", command=self.browse_save_location)
        self.browse_button.grid(row=2, column=2, padx=(0,10), pady=5, sticky="w")


        # 生成选项
        self.gen_word_var = tk.BooleanVar(value=True)
        self.gen_ppt_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(root, text="生成 Word 文档 (.docx)", variable=self.gen_word_var).grid(row=3, column=0, columnspan=3, padx=10, pady=5, sticky="w")
        ttk.Checkbutton(root, text="生成 PPT 演示文稿 (.pptx)", variable=self.gen_ppt_var).grid(row=4, column=0, columnspan=3, padx=10, pady=5, sticky="w")

        # 开始处理按钮
        self.process_button = ttk.Button(root, text="开始处理", command=self.start_processing_thread)
        self.process_button.grid(row=5, column=0, columnspan=3, padx=10, pady=10)

        # 状态与日志区域
        ttk.Label(root, text="状态与日志:").grid(row=6, column=0, padx=10, pady=5, sticky="w")
        self.status_text = scrolledtext.ScrolledText(root, wrap=tk.WORD, width=80, height=15, state='disabled')
        self.status_text.grid(row=7, column=0, columnspan=3, padx=10, pady=5, sticky="nsew")

        # 文件最终保存位置 (可以保留，也可以考虑移除，因为时间戳文件夹会在日志中显示)
        ttk.Label(root, text="时间戳子文件夹位置:").grid(row=8, column=0, padx=10, pady=5, sticky="w")
        self.save_location_label = ttk.Label(root, text="- 未开始 -", foreground="blue", wraplength=450) # wraplength
        self.save_location_label.grid(row=8, column=1, columnspan=2, padx=10, pady=5, sticky="w")

        # 使文本区域和输入框可以随窗口缩放
        root.grid_columnconfigure(1, weight=1)
        root.grid_rowconfigure(7, weight=1) # 日志区域行

        # 定期检查队列以更新UI
        self.root.after(100, self.process_status_queue)

    def browse_save_location(self):
        """打开文件夹选择对话框并更新路径"""
        directory = filedialog.askdirectory()
        if directory: # 如果用户选择了文件夹
            self.selected_save_path = directory
            self.save_path_entry.config(state='normal')
            self.save_path_entry.delete(0, tk.END)
            self.save_path_entry.insert(0, self.selected_save_path)
            self.save_path_entry.config(state='readonly')
            log_status(self.status_queue, f"选择的保存路径: {self.selected_save_path}")


    def update_status_text(self, message):
        """安全地更新状态文本框"""
        self.status_text.config(state='normal')
        self.status_text.insert(tk.END, message + "\n")
        self.status_text.see(tk.END) # 滚动到底部
        self.status_text.config(state='disabled')

    def process_status_queue(self):
        """处理状态队列中的消息"""
        try:
            while True:
                message = self.status_queue.get_nowait()
                self.update_status_text(message)
        except queue.Empty:
            pass # 队列为空，什么也不做
        self.root.after(100, self.process_status_queue) # 再次安排检查

    def _processing_task(self, article_url, doc_prefix, gen_word, gen_ppt, base_save_folder): # 添加 base_save_folder
        """实际执行处理任务的函数（在单独线程中运行）"""
        self.update_status_text("开始处理任务...")
        self.save_location_label.config(text="- 处理中... -")

        if not article_url or not doc_prefix: # 这个检查在start_processing_thread中已经做过，但保留无妨
            log_status(self.status_queue, "错误：URL和文档名称前缀不能为空。")
            self.process_button.config(state='normal') # 重新启用按钮
            self.save_location_label.config(text="- 输入错误 -")
            return

        if not base_save_folder: # 检查是否已选择保存路径
            log_status(self.status_queue, "错误：请先选择一个保存文件夹。")
            self.process_button.config(state='normal')
            self.save_location_label.config(text="- 未选择保存文件夹 -")
            return

        current_session_folder = create_timestamped_folder(self.status_queue, base_save_folder) # 传递路径
        if not current_session_folder:
            log_status(self.status_queue, "无法创建时间戳输出文件夹，任务中止。")
            self.process_button.config(state='normal')
            self.save_location_label.config(text="- 文件夹创建失败 -")
            return

        downloaded_images = download_images_from_url(article_url, current_session_folder, self.status_queue)

        if downloaded_images:
            if gen_word:
                generate_word_document(doc_prefix, downloaded_images, current_session_folder, self.status_queue)
            if gen_ppt:
                generate_ppt_presentation(doc_prefix, downloaded_images, current_session_folder, self.status_queue)
            log_status(self.status_queue, "所有选定文档创建完成！")
            self.save_location_label.config(text=current_session_folder) # 显示完整的时间戳路径
        else:
            log_status(self.status_queue, "没有下载到图片，无法生成文档。")
            self.save_location_label.config(text="- 未下载到图片 -")

        self.process_button.config(state='normal') # 任务完成，重新启用按钮
        log_status(self.status_queue, "-------------------- 处理结束 --------------------")


    def start_processing_thread(self):
        """启动处理任务的线程"""
        article_url = self.url_entry.get().strip()
        doc_prefix = self.prefix_entry.get().strip()
        gen_word = self.gen_word_var.get()
        gen_ppt = self.gen_ppt_var.get()

        if not self.selected_save_path: # 检查是否已选择保存路径
            messagebox.showerror("输入错误", "请先选择一个保存文件夹！")
            return

        if not article_url:
            messagebox.showerror("输入错误", "请输入微信文章 URL！")
            return
        if not doc_prefix:
            messagebox.showerror("输入错误", "请输入文档名称前缀！")
            return
        if not gen_word and not gen_ppt:
            messagebox.showwarning("选择错误", "请至少选择一种要生成的文档类型 (Word 或 PPT)！")
            return

        self.process_button.config(state='disabled') # 禁用按钮防止重复点击
        self.status_text.config(state='normal')
        self.status_text.delete(1.0, tk.END) # 清空之前的日志
        self.status_text.config(state='disabled')

        # 创建并启动线程
        thread = threading.Thread(target=self._processing_task,
                                  args=(article_url, doc_prefix, gen_word, gen_ppt, self.selected_save_path), # 传递选择的路径
                                  daemon=True) # 设置为守护线程，主程序退出时线程也退出
        thread.start()

# --- 主程序入口 ---
if __name__ == '__main__':
    # 检查Pillow是否安装，因为它是动态加载的依赖
    try:
        from PIL import Image
    except ImportError:
        print("错误：Pillow库未安装。请运行 'pip install Pillow' 来安装它。")
        # 如果在打包环境中，可能需要更复杂的处理或确保Pillow被包含
        # 对于简单的脚本运行，这里可以提示用户并退出
        if messagebox.askyesno("依赖缺失", "Pillow库未安装，它是本程序运行必需的。\n是否尝试现在安装 (需要pip)？"):
            try:
                import subprocess
                subprocess.check_call([sys.executable, "-m", "pip", "install", "Pillow"])
                messagebox.showinfo("安装成功", "Pillow已安装，请重新运行程序。")
            except Exception as e:
                messagebox.showerror("安装失败", f"自动安装Pillow失败: {e}\n请手动运行 'pip install Pillow'")
        import sys
        sys.exit(1)


    main_root = tk.Tk()
    app = WeixinToolApp(main_root)
    main_root.mainloop()