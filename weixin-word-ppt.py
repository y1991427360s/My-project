import requests
from bs4 import BeautifulSoup
import datetime
import os
from docx import Document
from docx.shared import Cm as word_Cm
from pptx import Presentation
from pptx.util import Cm as ppt_Cm
from PIL import Image # 新增导入

# --- 常量定义 ---
USER_AGENT = 'Mozilla/5.0 (iPhone; CPU iPhone OS 11_0 like Mac OS X) AppleWebKit/604.1.38 (KHTML, like Gecko) Version/11.0 Mobile/15A372 Safari/604.1'
BASE_DESKTOP_FOLDER_NAME = "weixin_images"
DEFAULT_IMAGE_EXTENSION = "jpg"


# --- 辅助函数 ---
def create_timestamped_folder():
    """
    在桌面创建主文件夹（如果不存在），并在其中创建一个带时间戳的子文件夹用于存放本次运行的文件。
    返回创建的子文件夹的绝对路径。
    """
    curr_time = datetime.datetime.now()
    timestamp_str = datetime.datetime.strftime(curr_time, '%Y%m%d%H%M')

    desktop_path = os.path.join(os.path.expanduser("~"), "Desktop")
    main_folder_path = os.path.join(desktop_path, BASE_DESKTOP_FOLDER_NAME)

    if not os.path.exists(main_folder_path):
        os.makedirs(main_folder_path) # 使用 makedirs 以防万一，虽然这里只创建一级

    session_folder_path = os.path.join(main_folder_path, timestamp_str)
    if not os.path.exists(session_folder_path):
        os.makedirs(session_folder_path)
    
    return session_folder_path


def download_images_from_url(url: str, save_folder: str) -> list[str]:
    """
    从给定的微信公众号URL下载图片到指定的文件夹。
    返回成功下载的图片文件的完整路径列表。
    """
    headers = {'user-agent': USER_AGENT}
    downloaded_image_paths = []
    
    try:
        response = requests.get(url=url, headers=headers)
        response.raise_for_status() # 如果请求失败则抛出HTTPError
        html_content = response.content.decode('utf-8', errors='ignore') # 指定utf-8并忽略解码错误
    except requests.exceptions.RequestException as e:
        print(f"请求URL失败: {url}, 错误: {e}")
        return downloaded_image_paths

    soup = BeautifulSoup(html_content, 'lxml')
    image_tags = soup.select('img')
    
    image_counter = 0
    for img_tag in image_tags:
        img_data_src = img_tag.get("data-src")
        if not img_data_src:
            # print("跳过一个没有data-src属性的图片标签")
            continue

        img_extension = img_tag.get("data-type", DEFAULT_IMAGE_EXTENSION)
        if not img_extension: # 以防万一data-type是空字符串
            img_extension = DEFAULT_IMAGE_EXTENSION
            
        img_filename = f"{image_counter}.{img_extension}"
        img_full_path = os.path.join(save_folder, img_filename)

        try:
            img_response = requests.get(url=img_data_src, headers=headers)
            img_response.raise_for_status()
            with open(img_full_path, 'wb') as f:
                f.write(img_response.content)
            downloaded_image_paths.append(img_full_path)
            image_counter += 1
        except requests.exceptions.RequestException as e:
            print(f"下载图片失败: {img_data_src}, 错误: {e}")
        except IOError as e:
            print(f"保存图片失败: {img_full_path}, 错误: {e}")
        except Exception as e:
            print(f"处理图片时发生未知错误: {img_data_src}, 错误: {e}")
            
    print(f"此次一共成功保存图片 {image_counter} 张到文件夹: {save_folder}")
    return downloaded_image_paths


def generate_word_document(file_name_prefix: str, image_paths: list[str], save_folder: str):
    """
    根据提供的图片路径列表生成Word文档。
    图片将适应页面宽度并保持宽高比。
    """
    if not image_paths:
        print("没有图片可用于生成Word文档。")
        return

    doc = Document()
    
    # 设置页面边距为0，使图片可以填充整个页面宽度
    for section in doc.sections:
        section.left_margin = word_Cm(0)
        section.right_margin = word_Cm(0)
        section.top_margin = word_Cm(0)
        section.bottom_margin = word_Cm(0)
    
    # 获取可用的页面宽度 (通常是减去边距后的宽度，这里边距为0)
    # A4纸张宽度约为21cm。默认页面大小可能是Letter。
    # 为了简单起见，我们直接使用一个常见的页面宽度值，或者从section获取。
    page_width_cm = doc.sections[0].page_width.cm

    for img_path in image_paths:
        try:
            # 添加图片，设置宽度为页面宽度，高度将自动按比例调整
            doc.add_picture(img_path, width=word_Cm(page_width_cm))
            # 每张图片后可以添加一个分页符，如果希望每张图片占一页
            # doc.add_page_break() 
        except Exception as e:
            print(f"无法将图片添加到Word文档: {img_path}, 错误: {e}")

    output_filename = f"{file_name_prefix}.docx"
    output_full_path = os.path.join(save_folder, output_filename)
    try:
        doc.save(output_full_path)
        print(f"Word文档已成功保存到: {output_full_path}")
    except Exception as e:
        print(f"保存Word文档失败: {output_full_path}, 错误: {e}")


def generate_ppt_presentation(file_name_prefix: str, image_paths: list[str], save_folder: str):
    """
    根据提供的图片路径列表生成PPT演示文稿。
    每张图片占据一页幻灯片，居中显示并尽可能填满幻灯片（保持宽高比）。
    """
    if not image_paths:
        print("没有图片可用于生成PPT。")
        return

    prs = Presentation()
    slide_width_cm = prs.slide_width.cm
    slide_height_cm = prs.slide_height.cm
    slide_aspect_ratio = slide_width_cm / slide_height_cm

    for img_path in image_paths:
        try:
            slide_layout = prs.slide_layouts[5] # 使用空白布局
            slide = prs.slides.add_slide(slide_layout)

            with Image.open(img_path) as img:
                img_width_px, img_height_px = img.size
            
            img_aspect_ratio = img_width_px / img_height_px

            if img_aspect_ratio > slide_aspect_ratio:
                # 图片比幻灯片更宽（或同样比例但需要以宽度为基准）
                pic_display_width_cm = slide_width_cm
                pic_display_height_cm = pic_display_width_cm / img_aspect_ratio
            else:
                # 图片比幻灯片更高（或同样比例但需要以高度为基准）
                pic_display_height_cm = slide_height_cm
                pic_display_width_cm = pic_display_height_cm * img_aspect_ratio
            
            left_cm = (slide_width_cm - pic_display_width_cm) / 2
            top_cm = (slide_height_cm - pic_display_height_cm) / 2
            
            slide.shapes.add_picture(
                img_path, 
                ppt_Cm(left_cm), 
                ppt_Cm(top_cm), 
                width=ppt_Cm(pic_display_width_cm), 
                height=ppt_Cm(pic_display_height_cm)
            )
        except Exception as e:
            print(f"无法将图片添加到PPT: {img_path}, 错误: {e}")

    output_filename = f"{file_name_prefix}.pptx"
    output_full_path = os.path.join(save_folder, output_filename)
    try:
        prs.save(output_full_path)
        print(f"PPT演示文稿已成功保存到: {output_full_path}")
    except Exception as e:
        print(f"保存PPT失败: {output_full_path}, 错误: {e}")


# --- 主程序逻辑 ---
if __name__ == '__main__':
    article_url = input("请输入微信公众号文章URL：")
    document_name_prefix = input("请设置文档名称前缀：")

    if not article_url or not document_name_prefix:
        print("URL和文档名称前缀不能为空。程序退出。")
    else:
        current_session_folder = create_timestamped_folder()
        print(f"文件将保存在: {current_session_folder}")

        downloaded_images = download_images_from_url(article_url, current_session_folder)

        if downloaded_images:
            print("正在生成Word文档...")
            generate_word_document(document_name_prefix, downloaded_images, current_session_folder)
            
            print("正在生成PPT演示文稿...")
            generate_ppt_presentation(document_name_prefix, downloaded_images, current_session_folder)
            
            print("所有文档创建完成！")
        else:
            print("没有下载到图片，无法生成文档。")